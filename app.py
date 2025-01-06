from flask import Flask, request, jsonify, send_file
from flask_swagger_ui import get_swaggerui_blueprint
from pptx import Presentation
from fpdf import FPDF
import os

app = Flask(__name__)
#swagger stuff
SWAGGER_URL = '/swagger'
API_URL = '/static/swagger.json'
swaggerui_blueprint = get_swaggerui_blueprint(SWAGGER_URL, API_URL, config={'app_name': "PPTX to PDF Converter"})
app.register_blueprint(swaggerui_blueprint, url_prefix=SWAGGER_URL)

# converting function for PPTX to PDF
@app.route('/convert', methods=['POST'])
def convert_pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']

    if not file.filename.endswith('.pptx'):
        return jsonify({"error": "File must be a .pptx file"}), 400

    input_file_path = os.path.join('uploads', file.filename)
    output_file_path = input_file_path.replace('.pptx', '.pdf')
    os.makedirs('uploads', exist_ok=True)
    file.save(input_file_path)

    # Convert PPTX to PDF
    try:
        pdf = FPDF()
        prs = Presentation(input_file_path)

        for slide in prs.slides:
            pdf.add_page()
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text
                        pdf.set_font("Arial", size=12)
                        pdf.multi_cell(0, 10, text)

        pdf.output(output_file_path)

        os.remove(input_file_path)

        return send_file(output_file_path, as_attachment=True, mimetype='application/pdf')
#error statement
    except Exception as e:
        return jsonify({"error": f"An error occurred: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(debug=True)